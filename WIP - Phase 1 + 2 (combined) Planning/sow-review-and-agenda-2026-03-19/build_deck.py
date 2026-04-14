#!/usr/bin/env python3
"""
Generate the Project Strong CEO Pitch Deck (14 slides)
Gold's Gym x StackHack Labs — March 31, 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design Tokens ──────────────────────────────────────────────────────────
GOLD = RGBColor(0xF2, 0xC2, 0x3E)
GOLD_LIGHT = RGBColor(0xFF, 0xF0, 0xB3)
GOLD_DIM = RGBColor(0xE6, 0xA8, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_2 = RGBColor(0x2D, 0x2D, 0x2D)
MID = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF7, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE2, 0xDC)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xD4, 0x88, 0x00)

FONT_DISPLAY = "Arial Black"  # Fallback for Bebas Neue
FONT_BODY = "Calibri"         # Fallback for DM Sans

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper Functions ───────────────────────────────────────────────────────

def add_dark_bg(slide):
    """Fill slide background with dark color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

def add_light_bg(slide):
    """Fill slide background with light color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=14, color=DARK, bold=False, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_name=FONT_BODY, font_size=14, color=DARK,
             bold=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None,
             border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def style_table_cell(cell, text, font_name=FONT_BODY, font_size=11, color=DARK,
                     bold=False, fill_color=None, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    cell.text_frame.word_wrap = True
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide)

# Large title
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "GOLD'S GYM  x  STACKHACK LABS",
            FONT_DISPLAY, 44, GOLD, True, PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "Build  vs.  Assemble  vs.  Buy",
            FONT_DISPLAY, 32, WHITE, False, PP_ALIGN.CENTER)

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.0), Inches(3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

# Date & context
add_textbox(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
            "March 31, 2026  |  Executive Strategy Session  |  Block 1, Day 1",
            FONT_BODY, 16, MID, False, PP_ALIGN.CENTER)

# Bottom tagline
add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
            "A strategic decision session, not a status update.",
            FONT_BODY, 14, RGBColor(0x88, 0x88, 0x88), False, PP_ALIGN.CENTER)

set_notes(slide, """MIKE OPENS:
"Thank you for having us. We've spent the last several weeks embedded in your organization -- interviewing 8 departments, surveying 41 franchisees, visiting 5 locations in person, and cataloging over 150 requirements."

"Today we're going to walk you through what we found, how Gold's compares to the market, and what we recommend as a path forward."

Set the tone: This is a partnership, not a vendor pitch. We're here to help Gold's make the right decisions.""")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2: PROJECT STRONG — THE VISION
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Section number
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(1), Inches(0.5),
            "01", FONT_DISPLAY, 18, GOLD, True)

# Title
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "PROJECT STRONG",
            FONT_DISPLAY, 48, GOLD, True, PP_ALIGN.LEFT)

# Subtitle
add_textbox(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
            "The unified franchise model and operating system for Gold's Gym of the Future",
            FONT_BODY, 18, WHITE, False)

# Vision bullets - left column
left_box = add_textbox(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.8), "",
                       FONT_BODY, 15, WHITE)
tf = left_box.text_frame
tf.paragraphs[0].text = ""
items = [
    "The hub of human performance in sports,\nfitness and wellness",
    "Category leadership in Strength, Performance,\nand Aspirational Fitness Culture",
    "30,000 sq ft strength-forward flagship model",
    "NOT another HVLP gym",
    "Premium positioning: $49.95 - $69.95+",
    "Lean operations: GM + AGM model",
]
for i, item in enumerate(items):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

# Callout box on right - top quote
callout = add_rect(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(1.8),
                   fill_color=DARK_2, border_color=GOLD, border_width=Pt(3))
add_textbox(slide, Inches(7.3), Inches(2.7), Inches(5), Inches(1.6),
            '"Your first time here is NOT your first\ntime doing fitness. This is the place\nyou graduate to. It is home to the athlete,\nthe fitness enthusiast, and the hybrid\ntraining person."',
            FONT_BODY, 15, GOLD_LIGHT, False)

# Second callout - not intimidating
callout2 = add_rect(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(1.2),
                    fill_color=DARK_2, border_color=GOLD_DIM, border_width=Pt(2))
add_textbox(slide, Inches(7.3), Inches(4.7), Inches(5), Inches(1.0),
            '"Not intimidating... it is home.\nBelonging to Gold\'s means something."',
            FONT_BODY, 15, GOLD_LIGHT, True)

# Bottom bar
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
            "Gold's Gym is entering its next era.",
            FONT_DISPLAY, 20, GOLD, True)

set_notes(slide, """MIKE — Opens with the vision:
- Gold's Gym is entering its next era
- Project Strong = the unified franchise model for the Gold's Gym of the Future
- "The hub of human performance in sports, fitness and wellness"
- Category leadership in Strength, Performance, and Aspirational Fitness Culture
- This is NOT another HVLP gym — "Your first time here is NOT your first time doing fitness. This is the place you graduate to."
- "Not intimidating... it is home to the athlete, the fitness enthusiast, the hybrid training person"
- "Belonging to Gold's means something."

MATT — Reinforces:
- 30,000 sq ft strength-forward flagship
- Premium pricing double HVLP ($49.95-$69.95+)
- Lean operations model (GM + AGM)
- Precision health is in Gold's DNA — "we facilitate the old school training and the new age approach all in one place"
- Every element of Project Strong requires unified technology""")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3: CORE AMENITIES
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "02", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
            "THE GOLD'S GYM OF THE FUTURE — CORE AMENITIES",
            FONT_DISPLAY, 28, DARK, True)

# Amenity cards - 2 columns x 4 rows
amenities = [
    ("MASSIVE FREE WEIGHT ZONES", "Expanded dumbbells, Olympic platforms,\nspecialty bars, strength machines"),
    ("PERFORMANCE TURF & FUNCTIONAL", "HYROX-style training, conditioning,\nathletic development"),
    ("SMALL GROUP PERFORMANCE STUDIO", "Strength & conditioning classes,\nathlete-inspired group training"),
    ("GOLD'S GYM POSING ROOM", "Heritage amenity for bodybuilders,\nphysique athletes, creators"),
    ("CONTENT-ENABLED SPACES", "Creator-friendly zones, social engagement\n(15% or less of floor space)"),
    ("RECOVERY & BIOHACKING SUITE", "Cold plunge, infrared sauna, red light,\ncompression — off days become on days"),
    ("PRECISION HEALTH & NUTRITION", "Supplements, peptides, blood testing,\nonsite nutrition, Gold's retail"),
    ("DIGITAL TRAINING", "App-delivered workouts, AI coaching,\nperformance tracking — 90% to 91%"),
]

for i, (title, desc) in enumerate(amenities):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.1)
    top = Inches(1.5 + row * 1.35)

    # Card background
    card = add_rect(slide, left, top, Inches(5.8), Inches(1.2),
                    fill_color=WHITE, border_color=BORDER)

    # Gold left accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(5), Inches(1.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), Inches(5.3), Inches(0.35),
                title, FONT_DISPLAY, 13, DARK, True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.5), Inches(5.3), Inches(0.65),
                desc, FONT_BODY, 11, MID)

# Membership tiers bar at bottom
tier_bg = add_rect(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.6),
                   fill_color=DARK, border_color=GOLD, border_width=Pt(2))
add_textbox(slide, Inches(1), Inches(6.75), Inches(3.5), Inches(0.4),
            "BASIC  |  $49.95+", FONT_DISPLAY, 14, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.8), Inches(6.75), Inches(4), Inches(0.4),
            "PERFORMANCE+  |  $69.95+", FONT_DISPLAY, 14, GOLD, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.2), Inches(6.75), Inches(3), Inches(0.4),
            "LOCAL ADD-ONS", FONT_DISPLAY, 14, MID, True, PP_ALIGN.CENTER)

set_notes(slide, """MIKE — Reference the Austin location as the model:
- Walk through each amenity category
- Emphasize this is NOT the old Gold's model — no basketball courts, no lap pools
- Limited old school group exercise — classes are performance and small-group based
- Recovery suite is tier-gated (Performance+ only) — creates upsell revenue
- KEY MIKE LINE: "Off days become on days at Gold's" — the recovery suite means members come in on rest days. This is a retention and revenue insight.
- Precision health: peptides, blood testing, supplements — "We are all about human optimization and precision health"
- "From 90% to 91% — we know that your performance matters"

INSERT AUSTIN LOCATION PHOTOS here — show the actual facility layout and amenity spaces as the reference model.

Key message: "Every single amenity here requires technology behind it. That's why this isn't just a design exercise — it's a technology architecture project."

MATT — Can add: "We'll come back to this on Slide 9 where we map each amenity to the specific tech it needs." """)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4: THE PROBLEM — CURRENT TECH STACK
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "03", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "THE PROBLEM",
            FONT_DISPLAY, 40, RED, True)
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "Tech Stack =/= Growth & Scale  |  Clubs Operating as Local Business Owners",
            FONT_BODY, 18, MID)

# Problem bullets
problems = [
    ("FRAGMENTED", "6-7 core vendors, ~20 apps across the ecosystem"),
    ("MEMBER MANAGEMENT FOCUSED", "The hub exists, but it's limited to basic operations"),
    ("NO CONTROL OF MEMBER EXPERIENCE", "Tech providers are driving strategy, not Gold's"),
    ("INCONSISTENT WEBSITES", "All franchisees on different websites — no brand consistency"),
    ("APP IS UTILITY, NOT EXPERIENTIAL", "Informational tool, not a service platform"),
    ("NOT BUILT FOR NATIONAL MARKETING", "Can't run unified campaigns or track ROI"),
    ("TECH PROVIDER = YOUR STRATEGY", "Instead of a strategy with tech that supports it"),
]

for i, (title, desc) in enumerate(problems):
    top = Inches(2.0 + i * 0.72)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), top + Inches(0.05), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD if i < 6 else RED
    circle.line.fill.background()
    # Number text
    add_textbox(slide, Inches(0.8), top + Inches(0.05), Inches(0.35), Inches(0.35),
                str(i + 1), FONT_DISPLAY, 12, DARK, True, PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Inches(1.4), top, Inches(3.5), Inches(0.35),
                title, FONT_DISPLAY, 14, WHITE, True)
    # Description
    add_textbox(slide, Inches(5), top + Inches(0.02), Inches(7.5), Inches(0.35),
                desc, FONT_BODY, 13, RGBColor(0xBB, 0xBB, 0xBB))

# Stats bar at bottom
stats_bg = add_rect(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7),
                    fill_color=DARK_2, border_color=GOLD_DIM, border_width=Pt(1))
stats = [
    ("2.5 / 5", "OVERALL SATISFACTION"),
    ("2.3 / 5", "ABC FITNESS"),
    ("2.0 / 5", "MOTIONSOFT"),
    ("4.5 / 5", "MOTION VIBE"),
]
for i, (val, label) in enumerate(stats):
    x = Inches(1.2 + i * 3)
    clr = GREEN if "4.5" in val else GOLD if "2.5" in val else RED
    add_textbox(slide, x, Inches(6.5), Inches(1.2), Inches(0.4),
                val, FONT_DISPLAY, 18, clr, True, PP_ALIGN.CENTER)
    add_textbox(slide, x - Inches(0.3), Inches(6.9), Inches(1.8), Inches(0.3),
                label, FONT_BODY, 9, MID, False, PP_ALIGN.CENTER)

set_notes(slide, """MIKE — Tells the story of what they heard:
- KEY FRAMING: "Tech Stack does NOT equal Growth & Scale. Clubs are operating as local business owners, not as a unified franchise brand."
- This is the thesis of the problem: the current tech doesn't support franchise growth — it keeps each club as an island.
- Walk through stakeholder interviews at high level — 8 departments
- "Your franchisees rated overall tech satisfaction at 2.5 out of 5. That's not a technology problem — that's a competitive positioning problem."
- ELEVATE BULLET 7: "Your tech provider IS your strategy vs. having a strategy with tech that supports it" — this is the core insight. Linger here.

MATT — Adds technical depth:
- Current System Scores: ABC at 2.3/5, Motionsoft at 2.0/5, HubSpot CRM at 2.5/5, Motion Vibe at 4.5/5
- Technology Fragmentation: core gym management, CRM, websites, mobile apps are all disconnected
- Critical Gaps: no brand-level controls, no real-time revenue visibility, no unified member database
- "The gap isn't about any single system being bad — it's that nothing talks to anything else. You have islands of data that can't work together."

Key stat: Franchisee tech satisfaction 2.5/5, Website is #1 pain point, 20-30+ fragmented apps, 0 SMS capability""")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5: THE ASSEMBLE STRATEGY
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "04", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
            "BUILD  vs.  ASSEMBLE  vs.  BUY",
            FONT_DISPLAY, 32, DARK, True)

# Three columns
col_data = [
    {
        "title": "BUILD",
        "color": RED,
        "bg": WHITE,
        "border": RED,
        "badge": None,
        "bullets": [
            "Custom platform from scratch",
            "Maximum control & differentiation",
            "Highest cost & longest timeline",
            "Requires large engineering team",
            "8/10 team investment",
        ],
        "verdict": "Aspirational — requires\nmore organizational readiness",
    },
    {
        "title": "ASSEMBLE",
        "color": GOLD,
        "bg": RGBColor(0xFF, 0xFB, 0xEB),
        "border": GOLD,
        "badge": "RECOMMENDED",
        "bullets": [
            "Own the core: experience + data",
            "Best-of-breed vendor components",
            "Moderate cost & timeline",
            "Vendors compete, not lock you in",
            "5/10 team investment",
        ],
        "verdict": "Matches Gold's DNA —\nasset-light, franchise-growth focused",
    },
    {
        "title": "BUY",
        "color": RED,
        "bg": WHITE,
        "border": RED,
        "badge": None,
        "bullets": [
            "Single off-the-shelf platform",
            "Fastest deployment",
            "Least control, vendor lock-in",
            "No vendor covers all needs",
            "2/10 team investment",
        ],
        "verdict": "Limiting — buying gaps\nand renting your future",
    },
]

for i, col in enumerate(col_data):
    left = Inches(0.6 + i * 4.2)
    top = Inches(1.5)
    w = Inches(3.9)
    h = Inches(5.2)
    bw = Pt(4) if col["badge"] else Pt(2)

    # Card
    card = add_rect(slide, left, top, w, h, fill_color=col["bg"],
                    border_color=col["border"], border_width=bw)

    # Badge
    if col["badge"]:
        badge = add_rect(slide, left + Inches(0.8), top - Inches(0.15), Inches(2.2), Inches(0.35),
                         fill_color=GOLD)
        add_textbox(slide, left + Inches(0.8), top - Inches(0.15), Inches(2.2), Inches(0.35),
                    col["badge"], FONT_DISPLAY, 12, DARK, True, PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.3), w - Inches(0.4), Inches(0.5),
                col["title"], FONT_DISPLAY, 24, col["color"], True, PP_ALIGN.CENTER)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 left + Inches(0.5), top + Inches(0.9),
                                 w - Inches(1), Pt(2))
    div.fill.solid()
    div.fill.fore_color.rgb = col["color"]
    div.line.fill.background()

    # Bullets
    for j, bullet in enumerate(col["bullets"]):
        bullet_y = top + Inches(1.1 + j * 0.55)
        # Bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     left + Inches(0.35), bullet_y + Inches(0.08),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = col["color"]
        dot.line.fill.background()
        add_textbox(slide, left + Inches(0.55), bullet_y, w - Inches(0.8), Inches(0.45),
                    bullet, FONT_BODY, 12, DARK)

    # Verdict
    verdict_y = top + h - Inches(1.1)
    div2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left + Inches(0.3), verdict_y - Inches(0.1),
                                  w - Inches(0.6), Pt(1))
    div2.fill.solid()
    div2.fill.fore_color.rgb = BORDER
    div2.line.fill.background()

    add_textbox(slide, left + Inches(0.3), verdict_y + Inches(0.05), w - Inches(0.6), Inches(0.8),
                col["verdict"], FONT_BODY, 12, col["color"], True, PP_ALIGN.CENTER)

set_notes(slide, """MIKE — Frames the decision simply:
- "There are really only three paths forward."
- BUILD: Custom platform, full control, highest investment. "Powerful model — but requires more organizational readiness."
- BUY: Single vendor, fastest but limiting. "No single vendor covers everything Gold's needs."
- ASSEMBLE (recommended): Best-in-class vendors integrated through a layer you own.
- IMPORTANT FRAMING: Don't position build as wrong — position it as aspirational. Assemble is the pragmatic starting point. "The architecture is the same either way."

MATT — Adds the org and revenue angle:
- Org requirements: Build 8/10, Assemble 5/10, Buy 2/10
- Revenue opportunity: "Gold's currently collects $0 in technology fees. Competitors charge $5+ per join, monthly tech fees."
- "Whether you assemble or build, this revenue opportunity is there."

KEY THREAD: "The architecture is the same regardless — owned data layer, owned member touchpoints. The question is how you get there." """)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6: THE 21ST CENTURY APPROACH — 7 PILLARS
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "05", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "AN APPROACH FOR THE 21ST CENTURY",
            FONT_DISPLAY, 32, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(8), Inches(0.4),
            "7 shifts from the current state to Project Strong",
            FONT_BODY, 16, MID)

# Table
tbl = add_table(slide, 8, 3, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.0))

# Header row
headers = ["", "FROM  (TODAY)", "TO  (PROJECT STRONG)"]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    style_table_cell(cell, h, FONT_DISPLAY, 12, GOLD, True, DARK_2, PP_ALIGN.CENTER)

# Column widths
tbl.columns[0].width = Inches(0.5)
tbl.columns[1].width = Inches(5.3)
tbl.columns[2].width = Inches(5.9)

rows_data = [
    ("1", "Fragmented data", "Data-first architecture"),
    ("2", "Vendor-controlled experience", "Own the member experience layer"),
    ("3", "Point-to-point integrations", "Vendors plug into your middleware"),
    ("4", "Franchisees on separate websites", "All franchisees on one corporate website"),
    ("5", "Separate utility apps", "One digital app for ALL offerings"),
    ("6", "App as info tool", "App is a 24/7 service platform for members"),
    ("7", "Manual operations", "Tech stack supported with AI"),
]

for i, (num, from_text, to_text) in enumerate(rows_data):
    r = i + 1
    bg = DARK if r % 2 == 0 else DARK_2
    style_table_cell(tbl.cell(r, 0), num, FONT_DISPLAY, 14, GOLD, True, bg, PP_ALIGN.CENTER)
    style_table_cell(tbl.cell(r, 1), from_text, FONT_BODY, 14, RGBColor(0xAA, 0xAA, 0xAA), False, bg)
    style_table_cell(tbl.cell(r, 2), to_text, FONT_BODY, 14, WHITE, True, bg)

set_notes(slide, """Present the recommended architecture as 7 clear shifts from current state.

Walk through each row — the "From" is where they are, the "To" is where Project Strong takes them.

Key emphasis points:
- Row 1 (Data-first): This is the foundation. Everything else depends on this.
- Row 3 (Middleware): Gold's controls the integration layer, vendors plug in.
- Row 4 (Website): Brand consistency — every franchisee on one website.
- Row 5-6 (App): The app goes from "check hours and freeze membership" to "your 24/7 gym experience."
- Row 7 (AI): "Brands that own their data today will deploy AI tomorrow."

This slide sets up the architecture slides that follow.""")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7: TECH STACK ARCHITECTURE — CAPABILITY MAP
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "06", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
            "TECH STACK ARCHITECTURE — CAPABILITY MAP",
            FONT_DISPLAY, 28, DARK, True)
add_textbox(slide, Inches(0.8), Inches(1.15), Inches(9), Inches(0.4),
            "What Gold's needs technology for — independent of any vendor",
            FONT_BODY, 15, MID)

# Placeholder rectangle for bubble chart
placeholder = add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8),
                       fill_color=LIGHT, border_color=BORDER, border_width=Pt(2))

# Zone labels inside placeholder
zones = [
    ("DATA & ANALYTICS", "Data Lake  |  Transformation  |  BI  |  Fran Connect", Inches(4.5), Inches(3.2), DARK),
    ("SALES & MARKETING", "CRM  |  CDP  |  Lead Mgmt  |  Marketing  |  Guest Pass", Inches(0.5), Inches(1.0), DARK_2),
    ("OPERATIONS", "Member Mgmt  |  POS  |  Scheduling  |  Staff  |  Gate  |  Inventory", Inches(0.5), Inches(5.0), DARK_2),
    ("MEMBER EXPERIENCE", "App  |  Content  |  Digital Training  |  Recovery Booking", Inches(8.0), Inches(3.0), DARK_2),
]

for zone_title, zone_items, zx, zy, zc in zones:
    # Zone bubble
    bubble = add_rect(slide, Inches(1.0) + zx, Inches(1.5) + zy, Inches(3.5), Inches(1.2),
                      fill_color=None, border_color=GOLD_DIM, border_width=Pt(2))
    add_textbox(slide, Inches(1.1) + zx, Inches(1.55) + zy, Inches(3.3), Inches(0.4),
                zone_title, FONT_DISPLAY, 13, zc, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.1) + zx, Inches(1.95) + zy, Inches(3.3), Inches(0.6),
                zone_items, FONT_BODY, 10, MID, False, PP_ALIGN.CENTER)

# Note
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
            "PLACEHOLDER: Replace with polished bubble chart from final report (environment-bubbles diagram). Do NOT name specific vendors.",
            FONT_BODY, 10, RGBColor(0xCC, 0x66, 0x00), True)

set_notes(slide, """MATT — Walk through the capability map:
- Data & Analytics at the center — the brain
- Sales & Marketing — CRM, CDP, lead management, marketing automation
- Operations — member management, POS, scheduling, gate control, inventory
- Member Experience — app, content, digital training, recovery booking

KEY POINT: "Do NOT name specific vendors on this slide. This is about what you NEED, not who provides it. That comes later."

This establishes the full scope of technology Gold's requires. The next slide shows which parts Gold's owns vs. which go to vendor RFP.

PRESENTATION NOTE (from Mike): These two slides (7 and 8) are intended as a BEFORE/AFTER REVEAL. Show this slide first as the plain capability map, then click to the next slide to reveal the color-coded version showing what Gold's owns vs. what goes to RFP. Build the anticipation: "Here's everything you need... now let's talk about what YOU own vs. what vendors provide." """)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8: OWN vs. VENDOR (RFI/RFP)
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "07", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
            "TECH STACK — OWN vs. VENDOR (RFI/RFP)",
            FONT_DISPLAY, 28, DARK, True)

# 3-layer architecture diagram
# Top layer - Gold's Owns (gold)
top_layer = add_rect(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.5),
                     fill_color=RGBColor(0xFF, 0xFB, 0xEB), border_color=GOLD, border_width=Pt(3))
add_textbox(slide, Inches(1.7), Inches(1.65), Inches(9.8), Inches(0.4),
            "GOLD'S OWNS — MEMBER-FACING TOUCHPOINTS",
            FONT_DISPLAY, 16, GOLD_DIM, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.7), Inches(2.1), Inches(9.8), Inches(0.8),
            "Mobile App  |  Member Portal  |  Franchise Dashboards  |  Corporate Website  |  Join Online",
            FONT_BODY, 14, DARK, False, PP_ALIGN.CENTER)

# Middle layer - Gold's Owns (gold, strongest)
mid_layer = add_rect(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.5),
                     fill_color=GOLD, border_color=GOLD_DIM, border_width=Pt(3))
add_textbox(slide, Inches(1.7), Inches(3.35), Inches(9.8), Inches(0.4),
            "GOLD'S OWNS — DATA & INTELLIGENCE LAYER",
            FONT_DISPLAY, 16, DARK, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.7), Inches(3.8), Inches(9.8), Inches(0.8),
            "Unified Data Model  |  Analytics & BI  |  AI / Automation  |  Open API  |  Middleware",
            FONT_BODY, 14, DARK_2, False, PP_ALIGN.CENTER)

# Bottom layer - Vendor (gray)
bot_layer = add_rect(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.5),
                     fill_color=LIGHT, border_color=BORDER, border_width=Pt(2))
add_textbox(slide, Inches(1.7), Inches(5.05), Inches(9.8), Inches(0.4),
            "VENDOR SYSTEMS OF RECORD  —  RFI / RFP",
            FONT_DISPLAY, 16, MID, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.7), Inches(5.5), Inches(9.8), Inches(0.8),
            "Core Gym Mgmt  |  CRM  |  Fitness/PT  |  Access Control  |  POS/Billing  |  Group Exercise",
            FONT_BODY, 14, MID, False, PP_ALIGN.CENTER)

# Arrows between layers
for y in [Inches(3.1), Inches(4.8)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), y, Inches(0.5), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD_DIM
    arrow.line.fill.background()

# Callout
add_textbox(slide, Inches(1.5), Inches(6.7), Inches(10), Inches(0.5),
            '"Own the brain (data + experience). Vendors provide the arms and legs."',
            FONT_BODY, 14, DARK, True, PP_ALIGN.CENTER)

set_notes(slide, """MATT — This is a KEY slide. Walk the 3-layer architecture:

TOP: Member-facing touchpoints — app, portal, dashboards, website — GOLD'S OWNS THESE.
"When a member opens the Gold's app, that should be YOUR experience, not ABC's, not a white-label."

MIDDLE: Gold's Data & Intelligence Layer — unified data model, analytics, AI, open API — GOLD'S OWNS THIS.
"This is the single most strategic decision you'll make. The data layer is where AI lives. Where automation lives. Where personalization lives. If you let vendors own this, you're renting your future."
"The middle layer is the moat. It's what lets you switch vendors without losing your data. It's what lets you deploy AI agents across all your systems."

BOTTOM: Vendor systems of record — they do the heavy lifting, you control the integration.

PLANT THE SEED: "Over time, as you build capability in this layer, you may find that building certain components yourself gives you an even bigger competitive edge. The architecture supports that evolution naturally." """)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9: TECHNOLOGY TIED TO AMENITIES
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "08", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
            "TECHNOLOGY TIED TO AMENITIES",
            FONT_DISPLAY, 28, DARK, True)
add_textbox(slide, Inches(0.8), Inches(1.15), Inches(9), Inches(0.4),
            "Every amenity in the gym needs technology behind it",
            FONT_BODY, 15, MID)

tbl = add_table(slide, 9, 2, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.2))
tbl.columns[0].width = Inches(4.2)
tbl.columns[1].width = Inches(7.5)

# Header
style_table_cell(tbl.cell(0, 0), "AMENITY", FONT_DISPLAY, 13, WHITE, True, DARK)
style_table_cell(tbl.cell(0, 1), "TECHNOLOGY REQUIRED", FONT_DISPLAY, 13, WHITE, True, DARK)

amenity_tech = [
    ("Free Weights / Strength Floor", "Gate access, member check-in, occupancy tracking"),
    ("Performance Turf / HYROX", "Class scheduling, booking, performance tracking"),
    ("Small Group Studio", "Class management, instructor scheduling, waitlists"),
    ("Recovery & Biohacking Suite", "Tier-based access control, booking, session tracking"),
    ("Gold's Gym Posing Room", "Reservation system, content/camera integration"),
    ("Content-Enabled Zones", "Creator tools, social integration, lighting controls"),
    ("Precision Health & Nutrition", "POS, inventory, member account linking, supplement/peptide tracking"),
    ("Digital Training", "App delivery, workout programming, AI coaching, performance data"),
]

for i, (amenity, tech) in enumerate(amenity_tech):
    r = i + 1
    bg = WHITE if r % 2 == 1 else LIGHT
    style_table_cell(tbl.cell(r, 0), amenity, FONT_DISPLAY, 12, DARK, True, bg)
    style_table_cell(tbl.cell(r, 1), tech, FONT_BODY, 12, DARK, False, bg)

set_notes(slide, """This slide connects the abstract tech conversation to the physical gym.

"Tech isn't abstract — every single amenity in your gym needs technology behind it."

Walk through each row:
- Free weights need gate access and occupancy tracking
- HYROX training needs class scheduling and performance tracking
- Recovery suite needs TIER-BASED access control (Performance+ members only) — this is revenue gating via technology
- Posing room needs reservation and content/camera integration
- Content zones need creator tools and social integration
- Nutrition/retail needs POS tied to member accounts
- Digital training is all app-delivered

Key message: "This is why you can't just have a gym management system. You need an ecosystem."

MATT: "And every one of these feeds data back into the middle layer we showed on the last slide." """)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10: MARKET CONTEXT & COMPETITIVE POSITIONING
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "09", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "MARKET CONTEXT & COMPETITIVE POSITIONING",
            FONT_DISPLAY, 28, GOLD, True)

# Stats bar
stats_bar = add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.7),
                     fill_color=DARK_2, border_color=GOLD_DIM, border_width=Pt(1))
stat_items = [("7", "BRANDS ANALYZED"), ("13", "TECH CATEGORIES"), ("91", "DATA POINTS")]
for i, (val, label) in enumerate(stat_items):
    x = Inches(1.5 + i * 4)
    add_textbox(slide, x, Inches(1.5), Inches(1), Inches(0.4),
                val, FONT_DISPLAY, 22, GOLD, True, PP_ALIGN.CENTER)
    add_textbox(slide, x - Inches(0.5), Inches(1.85), Inches(2), Inches(0.3),
                label, FONT_BODY, 10, MID, False, PP_ALIGN.CENTER)

# Premium competitors
add_textbox(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(0.4),
            "PREMIUM / PERFORMANCE BRANDS", FONT_DISPLAY, 14, WHITE, True)
premium = ["Equinox — App-first, data-driven, unified experience",
           "Life Time — Integrated platform, member ecosystem",
           "F45 — Franchise tech standard, unified ops",
           "Barry's — Booking-first, class-centric model",
           "Orangetheory — Wearable-integrated, performance data"]
pbox = add_textbox(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(2.5), "", FONT_BODY, 13, WHITE)
for i, item in enumerate(premium):
    p = pbox.text_frame.add_paragraph() if i > 0 else pbox.text_frame.paragraphs[0]
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.space_before = Pt(6)

# HVLP competitors
add_textbox(slide, Inches(7), Inches(2.6), Inches(5), Inches(0.4),
            "HVLP COMPETITORS", FONT_DISPLAY, 14, MID, True)
hvlp = ["Planet Fitness — Basic tech, volume play",
        "Crunch — Standard franchise ops",
        "EOS Fitness — Regional, limited tech"]
hbox = add_textbox(slide, Inches(7), Inches(3.0), Inches(5), Inches(1.5), "", FONT_BODY, 13, MID)
for i, item in enumerate(hvlp):
    p = hbox.text_frame.add_paragraph() if i > 0 else hbox.text_frame.paragraphs[0]
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    p.space_before = Pt(6)

# Key insight callout
insight_bg = add_rect(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5),
                      fill_color=DARK_2, border_color=GOLD, border_width=Pt(3))
add_textbox(slide, Inches(1.2), Inches(5.4), Inches(10.5), Inches(0.5),
            "THE OPPORTUNITY",
            FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(1.2), Inches(5.85), Inches(10.5), Inches(0.8),
            "None of the strength-focused brands have built this yet. Gold's can lead, not follow.\nThe assemble approach positions Gold's to leapfrog competitors without the cost of building from zero.",
            FONT_BODY, 15, WHITE)

set_notes(slide, """MIKE — Sets up the competitive frame:
- "We analyzed 7 brands across 13 technology categories — 91 data points."
- Every single one has invested in owning their member technology experience. This is table stakes.
- Point out where Gold's has gaps vs. competitors

MATT — Drives the strategic insight (KEY MOMENT):
- "Here's the most important finding: No one built a monolithic platform. Every franchisor assembles best-in-class components."
- "But the ones winning — they all own two things: the member-facing experience and the data layer underneath."
- "Planet Fitness doesn't build their own gym management software. But they absolutely own the member app experience and the data pipeline."
- Connect to AI: "The brands that own their data today are the ones that will deploy AI and automation tomorrow."
- "None of the strength-focused brands have built this yet. Gold's can be FIRST."

KEY NUANCE: "The most successful brands started by assembling — then over time began building custom components where it gave them a competitive edge." """)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11: NEXT STEPS — ARCHITECT THE PLAN
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "10", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
            "NEXT STEPS — ARCHITECT THE PLAN",
            FONT_DISPLAY, 28, DARK, True)

# Roadmap arrows
steps = [
    ("TODAY", "Phase 0 Complete\nDiscovery Done"),
    ("FEATURE SPECS", "Phase 1\nStrategy & Playbook"),
    ("VENDOR SELECTION", "Phase 2-3\nArchitecture & RFP"),
    ("BUILD & LAUNCH", "Phase 4-6\nImplement & Rollout"),
]

for i, (title, desc) in enumerate(steps):
    left = Inches(0.5 + i * 3.2)
    top = Inches(2.0)

    # Arrow shape
    if i < len(steps) - 1:
        arrow = add_rect(slide, left, top, Inches(3.0), Inches(1.8),
                         fill_color=GOLD if i == 0 else LIGHT,
                         border_color=GOLD if i == 0 else BORDER, border_width=Pt(2))
    else:
        arrow = add_rect(slide, left, top, Inches(3.0), Inches(1.8),
                         fill_color=LIGHT, border_color=BORDER, border_width=Pt(2))

    add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.6), Inches(0.4),
                title, FONT_DISPLAY, 16, DARK if i > 0 else WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.7), Inches(2.6), Inches(0.8),
                desc, FONT_BODY, 12, MID, False, PP_ALIGN.CENTER)

    # Connector arrow
    if i < len(steps) - 1:
        connector = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           left + Inches(3.0), top + Inches(0.65),
                                           Inches(0.3), Inches(0.5))
        connector.fill.solid()
        connector.fill.fore_color.rgb = GOLD
        connector.line.fill.background()

# Key message callout
callout_bg = add_rect(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.8),
                      fill_color=WHITE, border_color=GOLD, border_width=Pt(3))

# Gold left bar on callout
accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.5), Pt(6), Inches(1.8))
accent.fill.solid()
accent.fill.fore_color.rgb = GOLD
accent.line.fill.background()

add_textbox(slide, Inches(1.2), Inches(4.6), Inches(10.5), Inches(0.5),
            "YOU CAN'T WRITE AN RFP IF YOU DON'T KNOW WHAT YOU NEED",
            FONT_DISPLAY, 20, DARK, True)
add_textbox(slide, Inches(1.2), Inches(5.15), Inches(10.5), Inches(0.9),
            "Phase 1 delivers the detailed feature specs, vendor evaluation criteria, and architecture blueprints.\nThis is required to properly evaluate and select the right vendor partners.",
            FONT_BODY, 14, MID)

set_notes(slide, """MIKE — Walks the roadmap:
- "We need to architect the detailed plans and feature requirements."
- "This is required to properly evaluate and select the right vendor partners."
- "You can't write an RFP if you don't know what you need — that's what Phase 1 delivers."
- The output: detailed feature specs, vendor evaluation criteria, and architecture blueprints.

Keep this simple and forward-looking. The message is: "We know where we're going. Phase 1 maps the route."

Phase 1 & 2: $52,500 combined, 8-12 weeks
- Phase 1 (4-6 weeks): Strategy & Playbook
- Phase 2 (4-6 weeks): Architecture & Blueprint""")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12: TACTICS — USER STORY MAPPING
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "11", FONT_DISPLAY, 16, GOLD, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
            "TACTICS — HOW WE GET THERE",
            FONT_DISPLAY, 28, DARK, True)

# Timeline
timeline_steps = [
    ("THIS WEEK", "Executive Kickoff\n& Strategic Alignment"),
    ("WEEKS 2-4", "User Story\nMapping Sessions"),
    ("WEEKS 4-6", "Phase 1\nDeliverables"),
    ("WEEKS 6-8", "Phase 2\nArchitecture"),
]

for i, (when, what) in enumerate(timeline_steps):
    left = Inches(0.5 + i * 3.2)
    top = Inches(1.7)

    # Step circle/number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.2), top, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    add_textbox(slide, left + Inches(1.2), top + Inches(0.05), Inches(0.5), Inches(0.4),
                str(i + 1), FONT_DISPLAY, 16, DARK, True, PP_ALIGN.CENTER)

    # Connector line
    if i < len(timeline_steps) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      left + Inches(1.7), top + Inches(0.22),
                                      Inches(2.7), Pt(3))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD_LIGHT
        line.line.fill.background()

    # Labels
    add_textbox(slide, left + Inches(0.2), top + Inches(0.7), Inches(2.6), Inches(0.35),
                when, FONT_DISPLAY, 13, GOLD_DIM, True, PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), top + Inches(1.05), Inches(2.6), Inches(0.7),
                what, FONT_BODY, 13, DARK, False, PP_ALIGN.CENTER)

# Detail cards
cards = [
    ("USER STORY MAPPING SESSIONS",
     "Deep-dive sessions with domain experts — not just execs.\nGet into the weeds on features, workflows, and edge cases."),
    ("15 PERSONAS IDENTIFIED",
     "From Platform Admin to Prospective Member to Franchise Owner.\n100+ user stories ready for validation."),
    ("OUTPUT FEEDS DIRECTLY INTO DELIVERABLES",
     "Ensures we can tell vendors exactly what we need.\nNo ambiguity, no scope creep."),
]

for i, (title, desc) in enumerate(cards):
    left = Inches(0.5 + i * 4.2)
    top = Inches(4.0)

    card = add_rect(slide, left, top, Inches(3.9), Inches(2.5),
                    fill_color=WHITE, border_color=BORDER, border_width=Pt(1.5))

    # Gold top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.9), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(3.5), Inches(0.5),
                title, FONT_DISPLAY, 13, DARK, True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.9), Inches(3.5), Inches(1.3),
                desc, FONT_BODY, 12, MID)

set_notes(slide, """MIKE — Walks through the tactical plan:
- "User story mapping sessions with the Gold's team — not just execs, domain experts."
- "Get into the weeds on features, workflows, and edge cases."
- "These sessions happen in the weeks following this kickoff."
- "Output feeds directly into Phase 1 deliverables."
- "Ensures we can tell vendors EXACTLY what we need — no ambiguity, no scope creep."

Detail:
- 15 personas already identified (Platform Admin, Club Manager, Member, Franchise Owner, etc.)
- 100+ user stories drafted and ready for validation
- Sessions will involve: Brad, Danny, Robbie, Ryan, Rick + domain leads as needed

This is the bridge between strategy and execution.""")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13: KEY DECISIONS
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.4),
            "12", FONT_DISPLAY, 16, RED, True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
            "KEY DECISIONS — WE NEED YOUR DIRECTION",
            FONT_DISPLAY, 28, DARK, True)
add_textbox(slide, Inches(0.8), Inches(1.15), Inches(9), Inches(0.4),
            "Present each decision. State our recommendation. Ask for yes or no.",
            FONT_BODY, 14, MID, True)

tbl = add_table(slide, 6, 3, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tbl.columns[0].width = Inches(0.5)
tbl.columns[1].width = Inches(5.5)
tbl.columns[2].width = Inches(5.7)

# Header
style_table_cell(tbl.cell(0, 0), "#", FONT_DISPLAY, 12, WHITE, True, DARK, PP_ALIGN.CENTER)
style_table_cell(tbl.cell(0, 1), "DECISION", FONT_DISPLAY, 12, WHITE, True, DARK)
style_table_cell(tbl.cell(0, 2), "OUR RECOMMENDATION", FONT_DISPLAY, 12, GOLD, True, DARK)

decisions = [
    ("1",
     "Are we missing any core amenities to account for in the tech stack?",
     "Austin model is the baseline — confirm or flag additions"),
    ("2",
     "Should Gold's hire internally or engage an agency to build the owned components?",
     "Start with agency (speed + expertise), path to in-house as platform matures"),
    ("3",
     "Will this platform generate franchisee revenue (join fees, tech fees, paid services)?",
     "Yes — design revenue capture in from day one"),
    ("4",
     "Will the team support corporate only, or also franchisee services? More headcount needed?",
     "Shared services model — corporate team + dedicated franchisee support as platform rolls out"),
    ("5",
     "Is there a set budget, or do we provide tiered recommendations?",
     "We'll deliver tiered budget options (minimum, recommended, ideal) in Phase 1"),
]

for i, (num, decision, rec) in enumerate(decisions):
    r = i + 1
    bg = WHITE if r % 2 == 1 else LIGHT
    style_table_cell(tbl.cell(r, 0), num, FONT_DISPLAY, 16, GOLD, True, bg, PP_ALIGN.CENTER)
    style_table_cell(tbl.cell(r, 1), decision, FONT_BODY, 12, DARK, False, bg)
    style_table_cell(tbl.cell(r, 2), rec, FONT_BODY, 12, DARK, True, bg)

set_notes(slide, """THIS IS THE CEO DECISION SLIDE. This is where the 45-minute discussion happens.

Present each decision one at a time. State the recommendation. Ask for yes/no.

PRINT THIS SLIDE AS A SEPARATE HANDOUT for Brad and Danny.

Decision 1 — AMENITIES: "Are we missing anything? The Austin model is our baseline."
Decision 2 — BUILD TEAM: "Agency first for speed, path to in-house over time."
Decision 3 — REVENUE: "Yes — design revenue capture from day one. Join fees, tech fees, franchisee services."
Decision 4 — TEAM SCOPE: "Shared services model. Corporate team supports corporate locations AND provides franchisee services."
Decision 5 — BUDGET: "We'll provide tiered recommendations: minimum viable, recommended, and ideal."

MIKE leads this section. MATT supports with data/technical answers as questions arise.

Goal: Get directional yes/no on each so Phase 1 can proceed with clear guardrails.""")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14: CLOSE
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Large closing statement
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
            "GOLD'S GYM CREATED THE",
            FONT_DISPLAY, 36, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.6), Inches(11), Inches(1),
            "MODERN GYM INDUSTRY.",
            FONT_DISPLAY, 36, WHITE, True, PP_ALIGN.CENTER)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6), Inches(2.3), Pt(3))
div.fill.solid()
div.fill.fore_color.rgb = GOLD
div.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(1),
            "PROJECT STRONG WILL REDEFINE IT",
            FONT_DISPLAY, 36, GOLD, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(1),
            "FOR THE MODERN ERA.",
            FONT_DISPLAY, 36, GOLD, True, PP_ALIGN.CENTER)

# Tagline
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            "Thank you — let's build this.",
            FONT_BODY, 18, MID, False, PP_ALIGN.CENTER)

# Branding
add_textbox(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
            "Gold's Gym  x  StackHack Labs",
            FONT_DISPLAY, 16, RGBColor(0x55, 0x55, 0x55), False, PP_ALIGN.CENTER)

set_notes(slide, """MIKE — Closes with conviction:
"Gold's Gym created the modern gym industry. Project Strong will redefine it for the modern era."

"Thank you — let's build this."

This should feel like the natural conclusion. Confident, not salesy. The work speaks for itself.

After this slide, transition to the 45-minute decision discussion (Slide 13 decisions).
If time allows, offer to walk through any section of the interactive final report in more detail.""")


# ── Save ───────────────────────────────────────────────────────────────────
output_dir = "/Users/mattmasteller/Library/CloudStorage/OneDrive-SharedLibraries-StackHackLabsLLC/Gold's Gym One System - Documents/WIP - Phase 1 + 2 (combined) Planning/sow-review-and-agenda-2026-03-19"
output_path = os.path.join(output_dir, "Project-Strong-CEO-Pitch.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
